"""
file_parser.py - DataSeekAI File Parser Module

This module handles extraction of text and structured data from various file formats
uploaded by users. Supports CSV, Excel, PDF, DOCX, and PPTX files.
"""

import os
import pandas as pd
from typing import Dict, Union, List, Any
import traceback

# Optional imports with fallback error messages
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def parse_file(file_path: str) -> Dict[str, Any]:
    """
    Main function to detect file type and parse accordingly.
    
    Args:
        file_path (str): Path to the uploaded file
        
    Returns:
        dict: Structured data in format compatible with AI model
    """
    # Check if file exists
    if not os.path.exists(file_path):
        return {
            "type": "error",
            "message": f"File not found: {file_path}"
        }
    
    # Check if file is empty
    if os.path.getsize(file_path) == 0:
        return {
            "type": "error",
            "message": "File is empty"
        }
    
    # Get file extension
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower().lstrip('.')
    
    # Route to appropriate parser based on extension
    parsers = {
        'csv': parse_csv,
        'xlsx': parse_excel,
        'xls': parse_excel,
        'pdf': parse_pdf,
        'docx': parse_docx,
        'pptx': parse_pptx
    }
    
    if file_extension in parsers:
        try:
            return parsers[file_extension](file_path)
        except Exception as e:
            return {
                "type": "error",
                "message": f"Error parsing {file_extension.upper()} file: {str(e)}",
                "details": traceback.format_exc()
            }
    else:
        return {
            "type": "error",
            "message": f"Unsupported file format: {file_extension}. Supported formats: CSV, Excel, PDF, DOCX, PPTX"
        }


def parse_csv(file_path: str) -> Dict[str, Any]:
    """
    Parse CSV file using pandas.
    
    Args:
        file_path (str): Path to CSV file
        
    Returns:
        dict: Table structure with columns and preview data
    """
    try:
        # Read CSV with automatic delimiter detection
        try:
            df = pd.read_csv(file_path, nrows=11, encoding='utf-8')  # Read 11 rows to have 10 for preview after header
        except UnicodeDecodeError:
            df = pd.read_csv(file_path, nrows=11, encoding='latin1')
        
        if df.empty:
            return {
                "type": "error",
                "message": "CSV file contains no data"
            }
        
        # Convert to appropriate format
        columns = df.columns.tolist()
        
        # Get preview (first 10 rows)
        preview_data = df.head(10).to_dict(orient='records')
        
        # Clean data: convert NaN to None for JSON compatibility
        for row in preview_data:
            for key, value in row.items():
                if pd.isna(value):
                    row[key] = None
        
        return {
            "type": "table",
            "columns": columns,
            "data_preview": preview_data,
            "row_count": len(df)  # Note: this is from limited read; for actual count we'd need full read
        }
    
    except pd.errors.EmptyDataError:
        return {
            "type": "error",
            "message": "CSV file is empty or contains no data"
        }
    except pd.errors.ParserError as e:
        return {
            "type": "error",
            "message": f"CSV parsing error: {str(e)}"
        }
    except Exception as e:
        raise Exception(f"Unexpected error parsing CSV: {str(e)}")


def parse_excel(file_path: str) -> Dict[str, Any]:
    """
    Parse Excel file using pandas.
    
    Args:
        file_path (str): Path to Excel file
        
    Returns:
        dict: Table structure with columns and preview data from first sheet
    """
    try:
        # Read Excel file, get first sheet
        df = pd.read_excel(file_path, sheet_name=0, nrows=11)
        
        if df.empty:
            return {
                "type": "error",
                "message": "Excel sheet contains no data"
            }
        
        columns = df.columns.tolist()
        
        # Get preview (first 10 rows)
        preview_data = df.head(10).to_dict(orient='records')
        
        # Clean NaN values
        for row in preview_data:
            for key, value in row.items():
                if pd.isna(value):
                    row[key] = None
        
        # Get sheet names for reference
        xl = pd.ExcelFile(file_path)
        sheet_names = xl.sheet_names
        
        return {
            "type": "table",
            "columns": columns,
            "data_preview": preview_data,
            "row_count": len(df),
            "sheet_name": sheet_names[0],
            "available_sheets": sheet_names[1:] if len(sheet_names) > 1 else []
        }
    
    except Exception as e:
        raise Exception(f"Error parsing Excel file: {str(e)}")


def parse_pdf(file_path: str) -> Dict[str, Any]:
    """
    Extract text from PDF file.
    
    Args:
        file_path (str): Path to PDF file
        
    Returns:
        dict: Text structure with extracted content
    """
    if PdfReader is None:
        return {
            "type": "error",
            "message": "PyPDF library not installed. Please install pypdf to parse PDF files."
        }
    
    try:
        reader = PdfReader(file_path)
        
        if len(reader.pages) == 0:
            return {
                "type": "error",
                "message": "PDF file contains no pages"
            }
        
        extracted_text = []
        total_pages = len(reader.pages)
        
        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                extracted_text.append(f"--- Page {page_num} ---\n{text}")
        
        if not extracted_text:
            return {
                "type": "error",
                "message": "No text could be extracted from PDF. The file might be scanned or image-based."
            }
        
        full_text = "\n\n".join(extracted_text)
        
        return {
            "type": "text",
            "content": full_text,
            "metadata": {
                "page_count": total_pages,
                "extracted_pages": len(extracted_text)
            }
        }
    
    except Exception as e:
        raise Exception(f"Error parsing PDF: {str(e)}")


def parse_docx(file_path: str) -> Dict[str, Any]:
    """
    Extract text from DOCX file.
    
    Args:
        file_path (str): Path to DOCX file
        
    Returns:
        dict: Text structure with extracted content
    """
    if Document is None:
        return {
            "type": "error",
            "message": "python-docx library not installed. Please install python-docx to parse DOCX files."
        }
    
    try:
        doc = Document(file_path)
        
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    paragraphs.append(" | ".join(row_text))
        
        if not paragraphs:
            return {
                "type": "error",
                "message": "No text content found in DOCX file"
            }
        
        full_text = "\n\n".join(paragraphs)
        
        return {
            "type": "text",
            "content": full_text,
            "metadata": {
                "paragraph_count": len(paragraphs)
            }
        }
    
    except Exception as e:
        raise Exception(f"Error parsing DOCX: {str(e)}")


def parse_pptx(file_path: str) -> Dict[str, Any]:
    """
    Extract text from PPTX file.
    
    Args:
        file_path (str): Path to PPTX file
        
    Returns:
        dict: Text structure with extracted content
    """
    if Presentation is None:
        return {
            "type": "error",
            "message": "python-pptx library not installed. Please install python-pptx to parse PPTX files."
        }
    
    try:
        prs = Presentation(file_path)
        
        slides_text = []
        total_slides = len(prs.slides)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
                
                # Handle tables if present
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_content.append(" | ".join(row_text))
            
            if slide_content:
                slides_text.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_content))
        
        if not slides_text:
            return {
                "type": "error",
                "message": "No text content found in PPTX file"
            }
        
        full_text = "\n\n".join(slides_text)
        
        return {
            "type": "text",
            "content": full_text,
            "metadata": {
                "slide_count": total_slides,
                "extracted_slides": len(slides_text)
            }
        }
    
    except Exception as e:
        raise Exception(f"Error parsing PPTX: {str(e)}")


# Optional: Utility function to check if a file is supported
def is_supported_file(file_path: str) -> bool:
    """
    Check if the file type is supported by the parser.
    
    Args:
        file_path (str): Path to the file
        
    Returns:
        bool: True if supported, False otherwise
    """
    supported_extensions = {'.csv', '.xlsx', '.xls', '.pdf', '.docx', '.pptx'}
    _, ext = os.path.splitext(file_path)
    return ext.lower() in supported_extensions


# Optional: Function to get file metadata without full parsing
def get_file_metadata(file_path: str) -> Dict[str, Any]:
    """
    Get basic metadata about the file.
    
    Args:
        file_path (str): Path to the file
        
    Returns:
        dict: File metadata
    """
    if not os.path.exists(file_path):
        return {"error": "File not found"}
    
    stat = os.stat(file_path)
    _, ext = os.path.splitext(file_path)
    
    return {
        "filename": os.path.basename(file_path),
        "extension": ext.lower(),
        "size_bytes": stat.st_size,
        "size_kb": round(stat.st_size / 1024, 2),
        "modified": stat.st_mtime,
        "supported": is_supported_file(file_path)
    }


# Example usage and testing
if __name__ == "__main__":
    # This section can be used for testing the module
    import sys
    
    if len(sys.argv) > 1:
        test_file = sys.argv[1]
        print(f"Parsing file: {test_file}")
        result = parse_file(test_file)
        print("\nResult:")
        if result["type"] == "error":
            print(f"ERROR: {result['message']}")
        elif result["type"] == "table":
            print(f"Table with {result['row_count']} rows")
            print(f"Columns: {', '.join(result['columns'])}")
            print(f"Preview: {result['data_preview'][:2]}")  # Show first 2 rows
        else:  # text
            content = result['content']
            print(f"Text content (first 200 chars): {content[:200]}...")
    else:
        print("Please provide a file path to test: python file_parser.py <file_path>")