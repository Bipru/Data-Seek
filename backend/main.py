"""
main.py - DataSeekAI Backend Server

FastAPI-based backend for DataSeekAI that handles file uploads,
parsing, and AI-powered insights using Google Gemini.
"""

import os
import shutil
import uuid
from datetime import datetime
from typing import Optional, List, Dict, Any
from pathlib import Path

import pandas as pd
from fastapi import FastAPI, UploadFile, File, HTTPException, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import google.generativeai as genai
from dotenv import load_dotenv

# File parsing libraries (optional imports with graceful fallback)
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# Load environment variables
load_dotenv()

# Configure Gemini
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    print("WARNING: GEMINI_API_KEY not found in environment variables")
    print("Please set it in a .env file or environment")
else:
    genai.configure(api_key=GEMINI_API_KEY)

# Create FastAPI app
app = FastAPI(
    title="DataSeekAI API",
    description="Backend API for DataSeekAI chat interface",
    version="1.0.0"
)

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # In production, replace with specific origins
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Create uploads directory if it doesn't exist
UPLOAD_DIR = Path("uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

# Supported file types
SUPPORTED_EXTENSIONS = {'.csv', '.xlsx', '.xls', '.pdf', '.docx', '.pptx'}

# ==================== Pydantic Models ====================

class QueryRequest(BaseModel):
    """Request model for the /query endpoint"""
    user_message: str
    file_name: str

class QueryResponse(BaseModel):
    """Response model for the /query endpoint"""
    answer: str
    chart_type: str
    labels: List[str]
    values: List[float]
    scatter_data: Optional[List[Dict[str, Any]]] = None

class UploadResponse(BaseModel):
    """Response model for file upload"""
    filename: str
    file_size: int
    file_type: str
    message: str
    upload_id: str

class HealthResponse(BaseModel):
    """Health check response"""
    status: str
    timestamp: str
    version: str

# ==================== Helper Functions ====================

def get_file_extension(filename: str) -> str:
    """Extract and normalize file extension"""
    return Path(filename).suffix.lower()

def is_supported_file(filename: str) -> bool:
    """Check if file type is supported"""
    return get_file_extension(filename) in SUPPORTED_EXTENSIONS

def save_upload_file(upload_file: UploadFile) -> Path:
    """Save uploaded file to disk with unique name"""
    # Generate unique filename to avoid collisions
    original_name = Path(upload_file.filename)
    unique_name = f"{uuid.uuid4().hex}{original_name.suffix}"
    file_path = UPLOAD_DIR / unique_name
    
    # Save file
    with file_path.open("wb") as buffer:
        shutil.copyfileobj(upload_file.file, buffer)
    
    return file_path

def extract_file_content(file_path: Path) -> Dict[str, Any]:
    """
    Extract content from file based on its type.
    Returns a dictionary with type and extracted data.
    """
    extension = file_path.suffix.lower()
    
    if extension in ['.csv', '.xlsx', '.xls']:
        return extract_tabular_data(file_path, extension)
    elif extension == '.pdf':
        return extract_pdf_content(file_path)
    elif extension == '.docx':
        return extract_docx_content(file_path)
    elif extension == '.pptx':
        return extract_pptx_content(file_path)
    else:
        raise ValueError(f"Unsupported file type: {extension}")

def extract_tabular_data(file_path: Path, extension: str) -> Dict[str, Any]:
    """Extract data from CSV or Excel files"""
    try:
        if extension == '.csv':
            try:
                df = pd.read_csv(file_path, encoding='utf-8')
            except UnicodeDecodeError:
                df = pd.read_csv(file_path, encoding='latin1')
        else:  # Excel files
            df = pd.read_excel(file_path)
        
        # Get basic info
        columns = df.columns.tolist()
        
        # Get data preview (first 20 rows)
        preview = df.head(20).to_dict(orient='records')
        
        # Clean NaN values for JSON serialization
        for row in preview:
            for key, value in row.items():
                if pd.isna(value):
                    row[key] = None
        
        # Get data types
        dtypes = df.dtypes.astype(str).to_dict()
        
        # Get numeric columns for charting
        numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
        
        return {
            "type": "table",
            "columns": columns,
            "data_preview": preview,
            "row_count": len(df),
            "dtypes": dtypes,
            "numeric_columns": numeric_cols,
            "summary": df.describe().to_dict() if len(df) > 0 else {}
        }
    except Exception as e:
        raise Exception(f"Error parsing tabular file: {str(e)}")

def extract_pdf_content(file_path: Path) -> Dict[str, Any]:
    """Extract text from PDF file"""
    if PdfReader is None:
        raise ImportError("PyPDF library not installed")
    
    try:
        reader = PdfReader(file_path)
        text_content = []
        
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_content.append(text)
        
        full_text = "\n".join(text_content)
        
        return {
            "type": "text",
            "content": full_text,
            "page_count": len(reader.pages),
            "char_count": len(full_text)
        }
    except Exception as e:
        raise Exception(f"Error parsing PDF: {str(e)}")

def extract_docx_content(file_path: Path) -> Dict[str, Any]:
    """Extract text from DOCX file"""
    if Document is None:
        raise ImportError("python-docx library not installed")
    
    try:
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        
        full_text = "\n".join(paragraphs)
        
        return {
            "type": "text",
            "content": full_text,
            "paragraph_count": len(paragraphs),
            "char_count": len(full_text)
        }
    except Exception as e:
        raise Exception(f"Error parsing DOCX: {str(e)}")

def extract_pptx_content(file_path: Path) -> Dict[str, Any]:
    """Extract text from PPTX file"""
    if Presentation is None:
        raise ImportError("python-pptx library not installed")
    
    try:
        prs = Presentation(file_path)
        slides_text = []
        
        for slide in prs.slides:
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_content.append(shape.text)
            if slide_content:
                slides_text.extend(slide_content)
        
        full_text = "\n".join(slides_text)
        
        return {
            "type": "text",
            "content": full_text,
            "slide_count": len(prs.slides),
            "char_count": len(full_text)
        }
    except Exception as e:
        raise Exception(f"Error parsing PPTX: {str(e)}")

def prepare_ai_context(file_content: Dict[str, Any]) -> str:
    """Prepare context string for AI based on file content"""
    if file_content["type"] == "table":
        context = f"""
Dataset Information:
- Columns: {', '.join(file_content['columns'])}
- Total rows: {file_content['row_count']}
- Numeric columns: {', '.join(file_content.get('numeric_columns', []))}

Data Preview (first 5 rows):
{file_content['data_preview'][:5]}

Summary Statistics:
{file_content.get('summary', {})}
"""
    else:
        # For text documents, include excerpt
        content = file_content['content']
        excerpt = content[:2000] + "..." if len(content) > 2000 else content
        context = f"""
Document Content:
{excerpt}

Total characters: {file_content.get('char_count', 0)}
"""
    return context

async def query_gemini(context: str, user_question: str) -> Dict[str, Any]:
    """Query Gemini API for insights and chart recommendations"""
    if not GEMINI_API_KEY:
        # Fallback response when API key is missing
        return {
            "text": "Gemini API key not configured. Please set GEMINI_API_KEY in environment variables.",
            "chart_type": "none",
            "labels": [],
            "values": []
        }
    
    try:
        # Initialize Gemini model
        model = genai.GenerativeModel('gemini-2.5-flash')
        
        # Construct prompt
        prompt = f"""
You are DataSeekAI, an AI assistant that helps users analyze their data.

Context from uploaded file:
{context}

User question: {user_question}

Based on the context and question, provide:
1. A clear, concise insight answering the question
2. A recommended chart type if visualization would help (choose from: bar, line, pie, table, or none)
3. If a chart is recommended, provide appropriate labels and values

Format your response as valid JSON with these fields:
- "text": your insight explanation
- "chart_type": one of "bar", "line", "pie", "table", or "none"
- "labels": array of strings for chart labels
- "values": array of numbers for chart values

Only respond with the JSON, no additional text.
"""
        
        # Generate response
        response = model.generate_content(prompt)
        
        # Parse JSON from response
        try:
            # Extract JSON from response text
            response_text = response.text.strip()
            # Find JSON boundaries if needed
            if response_text.startswith('```json'):
                response_text = response_text[7:]
            if response_text.startswith('```'):
                response_text = response_text[3:]
            if response_text.endswith('```'):
                response_text = response_text[:-3]
            
            result = eval(response_text)  # Use eval for simplicity, but consider json.loads in production
            
            # Ensure required fields
            return {
                "text": result.get("text", "No insight generated"),
                "chart_type": result.get("chart_type", "none"),
                "labels": result.get("labels", []),
                "values": result.get("values", [])
            }
        except Exception as e:
            print(f"Error parsing Gemini response: {e}")
            print(f"Raw response: {response.text}")
            return {
                "text": "Error processing AI response. Please try again.",
                "chart_type": "none",
                "labels": [],
                "values": []
            }
            
    except Exception as e:
        print(f"Gemini API error: {str(e)}")
        return {
            "text": f"AI service error: {str(e)}",
            "chart_type": "none",
            "labels": [],
            "values": []
        }

# ==================== API Endpoints ====================

@app.get("/health", response_model=HealthResponse)
async def health_check():
    """Health check endpoint"""
    return {
        "status": "DataSeekAI backend running",
        "timestamp": datetime.now().isoformat(),
        "version": app.version
    }

@app.post("/upload", response_model=UploadResponse)
async def upload_file(
    file: UploadFile = File(...),
    background_tasks: BackgroundTasks = None
):
    """
    Upload and save a file for analysis.
    Supports: CSV, Excel, PDF, DOCX, PPTX
    """
    # Validate file type
    if not is_supported_file(file.filename):
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type. Supported types: {', '.join(SUPPORTED_EXTENSIONS)}"
        )
    
    try:
        # Save file
        file_path = save_upload_file(file)
        
        # Get file size
        file_size = file_path.stat().st_size
        
        # Return success response
        return UploadResponse(
            filename=file.filename,
            file_size=file_size,
            file_type=get_file_extension(file.filename),
            message="File uploaded successfully",
            upload_id=file_path.stem
        )
    
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Error saving file: {str(e)}"
        )

@app.post("/query", response_model=QueryResponse)
async def query_data(request: QueryRequest):
    """
    Query the AI about an uploaded file.
    Returns insights and chart recommendations.
    """
    # Validate input
    if not request.user_message.strip():
        raise HTTPException(
            status_code=400,
            detail="User message cannot be empty"
        )
    
    # Find the uploaded file
    try:
        # Look for file in uploads directory (simplified - in production, track uploads in DB)
        upload_files = list(UPLOAD_DIR.glob(f"*{Path(request.file_name).suffix}"))
        
        if not upload_files:
            # Try to find by original name pattern
            upload_files = list(UPLOAD_DIR.glob("*"))
        
        if not upload_files:
            raise HTTPException(
                status_code=404,
                detail=f"File '{request.file_name}' not found. Please upload it first."
            )
        
        # Use the most recent matching file (simplified approach)
        file_path = max(upload_files, key=lambda p: p.stat().st_mtime)
        
    except Exception as e:
        raise HTTPException(
            status_code=404,
            detail=f"Error locating file: {str(e)}"
        )
    
    try:
        # Extract content from file
        file_content = extract_file_content(file_path)
        
        # Prepare context for AI
        context = prepare_ai_context(file_content)
        
        # Query Gemini
        ai_response = await query_gemini(context, request.user_message)
        
        # Return formatted response
        return QueryResponse(
            answer=ai_response["text"],
            chart_type=ai_response["chart_type"],
            labels=ai_response["labels"],
            values=ai_response["values"],
            scatter_data=ai_response.get("scatter_data", None)
        )
        
    except ImportError as e:
        raise HTTPException(
            status_code=500,
            detail=f"Missing required library: {str(e)}"
        )
    except ValueError as e:
        raise HTTPException(
            status_code=400,
            detail=str(e)
        )
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Error processing query: {str(e)}"
        )

@app.get("/files")
async def list_uploaded_files():
    """List all uploaded files (useful for debugging)"""
    files = []
    for file_path in UPLOAD_DIR.iterdir():
        if file_path.is_file():
            files.append({
                "filename": file_path.name,
                "size": file_path.stat().st_size,
                "modified": datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
                "extension": file_path.suffix
            })
    return {"files": files}

@app.delete("/files/{filename}")
async def delete_file(filename: str):
    """Delete an uploaded file"""
    file_path = UPLOAD_DIR / filename
    if file_path.exists():
        file_path.unlink()
        return {"message": f"File {filename} deleted"}
    raise HTTPException(status_code=404, detail="File not found")

# ==================== Startup/Shutdown Events ====================

@app.on_event("startup")
async def startup_event():
    """Run on application startup"""
    print(f"DataSeekAI backend starting up...")
    print(f"Upload directory: {UPLOAD_DIR.absolute()}")
    if not GEMINI_API_KEY:
        print("WARNING: GEMINI_API_KEY not configured")
    else:
        print("Gemini API configured")

@app.on_event("shutdown")
async def shutdown_event():
    """Run on application shutdown"""
    print("DataSeekAI backend shutting down...")

# ==================== Run with Uvicorn ====================

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(
        "main:app",
        host="0.0.0.0",
        port=8000,
        reload=True,  # Enable for development
        log_level="info"
    )